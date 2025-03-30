#!/usr/bin/env python3
# RAG와 에이전트에 대한 PowerPoint 프레젠테이션 생성 스크립트

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# 이미지 폴더 경로
IMAGE_DIR = "/Users/macbook/gdrive/SelfStudy/ds4th_study/source/AI_Engineering/images"

def create_presentation():
    # 새 프레젠테이션 생성
    prs = Presentation()
    
    # 기본 폰트 설정
    font_name = '맑은 고딕'
    
    # 제목 슬라이드 추가
    title_slide_layout = prs.slide_layouts[0]  # 제목 슬라이드 레이아웃
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "RAG와 에이전트: AI의 능력을 확장하는 두 가지 핵심 패턴"
    subtitle.text = "컨텍스트와 도구를 통한 모델 성능 향상"
    
    # 폰트 설정
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(40)
            run.font.bold = True
    
    for paragraph in subtitle.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(28)
    
    # 슬라이드 내용 정의
    slide_contents = [
        # 슬라이드 1: 소개
        {
            "title": "소개",
            "bullets": [
                "본 발표의 목적과 범위", 
                "다룰 주요 주제: RAG, 에이전트, 메모리 시스템", 
                "발표 시간 및 구성 안내"
            ],
            "image": None
        },
        # 슬라이드 2: AI 모델의 한계
        {
            "title": "AI 모델의 한계",
            "bullets": [
                "컨텍스트 크기 제한", 
                "학습 데이터의 시간적 제약 (오래된 정보)", 
                "환각(hallucination) 문제", 
                "한정된 기능 (텍스트 생성만 가능)"
            ],
            "image": None
        },
        # 슬라이드 3: 한계 극복을 위한 두 패턴
        {
            "title": "한계 극복을 위한 두 패턴",
            "bullets": [
                "RAG: 외부 지식을 통한 모델 보강", 
                "에이전트: 도구를 통한 모델 능력 확장", 
                "두 패턴의 공통점과 차이점", 
                "모델 자체 수정 없이 성능 향상 가능성"
            ],
            "image": None
        },
        # 슬라이드 4: RAG 개요
        {
            "title": "RAG 개요",
            "bullets": [
                "RAG의 정의: \"Retrieval-Augmented Generation\"", 
                "RAG의 기본 아이디어: 검색 + 생성", 
                "RAG가 해결하는 문제들", 
                "RAG의 주요 장점: 최신 정보 제공, 환각 감소, 맥락 인식 향상"
            ],
            "image": None
        },
        # 슬라이드 5: RAG 아키텍처
        {
            "title": "RAG 아키텍처",
            "bullets": [
                "fig_06_01 활용 (RAG 아키텍처 다이어그램)", 
                "검색기(Retriever)와 생성기(Generator) 설명", 
                "인덱스와 쿼리 처리의 중요성", 
                "전체 워크플로우 개요"
            ],
            "image": "fig_06_01.png"
        },
        # 슬라이드 6: RAG 데이터 처리
        {
            "title": "RAG 데이터 처리",
            "bullets": [
                "문서 수집 및 전처리", 
                "청킹(Chunking) 전략", 
                "메타데이터 활용 방법", 
                "효과적인 데이터 준비의 중요성"
            ],
            "image": None
        },
        # 슬라이드 7: 청킹 전략
        {
            "title": "청킹 전략",
            "bullets": [
                "fig_06_02 활용 (청킹 방법 비교)", 
                "다양한 청킹 방법: 고정 크기, 문단, 의미 기반", 
                "각 방법의 장단점", 
                "청킹이 검색 성능에 미치는 영향"
            ],
            "image": "fig_06_02.png"
        },
        # 슬라이드 8: 검색 알고리즘 I - 용어 기반 검색
        {
            "title": "검색 알고리즘 I - 용어 기반 검색",
            "bullets": [
                "BM25, TF-IDF 등 전통적 검색 알고리즘", 
                "키워드 매칭 기반의 장단점", 
                "구현 용이성과 계산 효율성", 
                "언제 용어 기반 검색이 적합한가?"
            ],
            "image": None
        },
        # 슬라이드 9: 검색 알고리즘 II - 임베딩 기반 검색
        {
            "title": "검색 알고리즘 II - 임베딩 기반 검색",
            "bullets": [
                "fig_06_05 활용 (임베딩 기반 검색 설명)", 
                "벡터 표현과 의미적 유사성", 
                "임베딩 모델 종류와 특성", 
                "계산 비용과 성능 비교"
            ],
            "image": "fig_06_05.png"
        },
        # 슬라이드 10: 검색 알고리즘 III - 하이브리드 접근법
        {
            "title": "검색 알고리즘 III - 하이브리드 접근법",
            "bullets": [
                "용어 기반 + 임베딩 기반 결합의 이점", 
                "검색 결과 재순위화(reranking) 전략", 
                "fig_06_06 활용 (하이브리드 검색 설명)", 
                "실제 사례와 성능 향상 예시"
            ],
            "image": "fig_06_06.png"
        },
        # 슬라이드 11: 쿼리 최적화
        {
            "title": "쿼리 최적화",
            "bullets": [
                "쿼리 재작성(Query rewriting) 기법", 
                "히포크래틱(Hypocratic) 검색", 
                "다중 쿼리 생성", 
                "검색 성능 향상을 위한 프롬프트 엔지니어링"
            ],
            "image": None
        },
        # 슬라이드 12: RAG 구현 고려사항
        {
            "title": "RAG 구현 고려사항",
            "bullets": [
                "임베딩 모델 선택", 
                "벡터 데이터베이스 비교", 
                "인덱싱 전략", 
                "구현 복잡성과 성능 트레이드오프"
            ],
            "image": None
        },
        # 슬라이드 13: RAG 평가
        {
            "title": "RAG 평가",
            "bullets": [
                "컨텍스트 정밀도와 재현율", 
                "검색 품질 평가 방법", 
                "응답 품질 평가 메트릭", 
                "평가 자동화 전략"
            ],
            "image": None
        },
        # 슬라이드 14: RAG 확장 - 멀티모달 RAG
        {
            "title": "RAG 확장 - 멀티모달 RAG",
            "bullets": [
                "텍스트를 넘어선 RAG 시스템", 
                "이미지, 오디오, 비디오 처리", 
                "멀티모달 임베딩 활용", 
                "실제 응용 사례"
            ],
            "image": None
        },
        # 슬라이드 15: RAG 확장 - 표 형식 데이터와 SQL
        {
            "title": "RAG 확장 - 표 형식 데이터와 SQL",
            "bullets": [
                "표 형식 데이터와 RAG 통합", 
                "SQL 쿼리 생성 및 실행", 
                "Kitty Vogue 예시 활용", 
                "엔터프라이즈 데이터 활용 시나리오"
            ],
            "image": None
        },
        # 슬라이드 16: 에이전트 개요
        {
            "title": "에이전트 개요",
            "bullets": [
                "에이전트의 정의", 
                "fig_06_08 활용 (에이전트 시각화)", 
                "에이전트가 가능하게 하는 새로운 애플리케이션", 
                "환경과 행동 집합의 중요성"
            ],
            "image": "fig_06_08.png"
        },
        # 슬라이드 17: 에이전트의 특성
        {
            "title": "에이전트의 특성",
            "bullets": [
                "환경(Environment) 인식 능력", 
                "행동 집합(Action set)과 도구(Tools)", 
                "자율성과 목표 지향성", 
                "계획 및 의사결정 능력"
            ],
            "image": None
        },
        # 슬라이드 18: 도구(Tools) I - 지식 보강
        {
            "title": "도구(Tools) I - 지식 보강",
            "bullets": [
                "텍스트 검색기, 이미지 검색기", 
                "웹 브라우징과 인터넷 API", 
                "내부 정보 접근", 
                "RAG를 에이전트의 도구로 활용"
            ],
            "image": None
        },
        # 슬라이드 19: 도구(Tools) II - 기능 확장
        {
            "title": "도구(Tools) II - 기능 확장",
            "bullets": [
                "계산기, 단위 변환기, 달력 등 간단한 도구", 
                "코드 해석기와 실행기", 
                "텍스트-이미지 생성 도구", 
                "멀티모달 확장 전략"
            ],
            "image": None
        },
        # 슬라이드 20: 도구(Tools) III - 쓰기 작업
        {
            "title": "도구(Tools) III - 쓰기 작업",
            "bullets": [
                "환경에 영향을 미치는 행동", 
                "이메일 API, SQL 실행기, 은행 API 등", 
                "보안 및 안전 고려사항", 
                "현실 세계 행동의 위험과 기회"
            ],
            "image": None
        },
        # 슬라이드 21: 도구 선택 전략
        {
            "title": "도구 선택 전략",
            "bullets": [
                "fig_06_14 활용 (도구 사용 패턴)", 
                "최적의 도구 세트 결정 방법", 
                "도구 복잡성과 유용성 균형", 
                "도구 사용 분석 및 최적화"
            ],
            "image": "fig_06_14.png"
        },
        # 슬라이드 22: 계획(Planning) 개요
        {
            "title": "계획(Planning) 개요",
            "bullets": [
                "계획의 정의와 중요성", 
                "계획과 실행의 분리 필요성", 
                "fig_06_09 활용 (계획 프로세스)", 
                "계획 생성, 검증, 실행의 3단계"
            ],
            "image": "fig_06_09.png"
        },
        # 슬라이드 23: 계획 생성
        {
            "title": "계획 생성",
            "bullets": [
                "모델을 계획 생성기로 전환하는 방법", 
                "프롬프트 엔지니어링 접근법", 
                "복잡한 작업 분해", 
                "계획의 세분화 수준(granularity)"
            ],
            "image": None
        },
        # 슬라이드 24: 함수 호출(Function Calling)
        {
            "title": "함수 호출(Function Calling)",
            "bullets": [
                "함수 호출의 개념과 구현", 
                "fig_06_10 활용 (함수 호출 예시)", 
                "도구 인벤토리 관리", 
                "API 매개변수 처리 전략"
            ],
            "image": "fig_06_10.png"
        },
        # 슬라이드 25: 제어 흐름(Control Flow)
        {
            "title": "제어 흐름(Control Flow)",
            "bullets": [
                "fig_06_11 활용 (제어 흐름 시각화)", 
                "순차, 병렬, if문, 루프 등 다양한 제어 흐름", 
                "복잡한 계획의 실행 관리", 
                "병렬 실행의 장점과 구현 방법"
            ],
            "image": "fig_06_11.png"
        },
        # 슬라이드 26: 반성(Reflection) 및 오류 수정
        {
            "title": "반성(Reflection) 및 오류 수정",
            "bullets": [
                "반성의 중요성과 타이밍", 
                "fig_06_12 활용 (ReAct 프레임워크)", 
                "반성 구현 방법: 자기비판, 평가자 활용", 
                "오류 감지 및 수정 메커니즘"
            ],
            "image": "fig_06_12.png"
        },
        # 슬라이드 27: 반성 프레임워크 - ReAct와 Reflexion
        {
            "title": "반성 프레임워크 - ReAct와 Reflexion",
            "bullets": [
                "ReAct: 추론과 행동 교차 수행", 
                "fig_06_13 활용 (Reflexion 에이전트)", 
                "계획 수정 및 개선 과정", 
                "비용과 지연 시간 트레이드오프"
            ],
            "image": "fig_06_13.png"
        },
        # 슬라이드 28: 에이전트 실패 모드
        {
            "title": "에이전트 실패 모드",
            "bullets": [
                "계획 실패: 잘못된 도구, 매개변수 오류 등", 
                "도구 실패: 도구 출력 오류, 번역 오류 등", 
                "효율성 문제: 불필요한 단계, 과도한 비용 등", 
                "실패 감지 및 디버깅 방법"
            ],
            "image": None
        },
        # 슬라이드 29: 에이전트 평가
        {
            "title": "에이전트 평가",
            "bullets": [
                "계획 품질 평가 지표", 
                "도구 사용 분석", 
                "작업 완료율 측정", 
                "비용 및 시간 효율성 평가"
            ],
            "image": None
        },
        # 슬라이드 30: 메모리 시스템 개요
        {
            "title": "메모리 시스템 개요",
            "bullets": [
                "fig_06_16 활용 (메모리 메커니즘)", 
                "내부 지식(Internal knowledge)", 
                "단기 메모리(Short-term memory)", 
                "장기 메모리(Long-term memory)"
            ],
            "image": "fig_06_16.png"
        },
        # 슬라이드 31: 메모리 관리 전략
        {
            "title": "메모리 관리 전략",
            "bullets": [
                "FIFO, 중복 제거, 요약 기반 접근법", 
                "모순 처리 방법", 
                "메모리 시스템이 에이전트에 미치는 영향", 
                "구현 사례 및 성능 영향"
            ],
            "image": None
        },
        # 슬라이드 32: RAG vs 에이전트 비교
        {
            "title": "RAG vs 에이전트 비교",
            "bullets": [
                "두 패턴의 핵심 차이점 요약", 
                "사용 사례별 적합성", 
                "결합 시 시너지 효과", 
                "선택 시 고려사항"
            ],
            "image": None
        },
        # 슬라이드 33: 실제 구현 사례
        {
            "title": "실제 구현 사례",
            "bullets": [
                "오픈소스 RAG 프레임워크", 
                "에이전트 구현 도구와 플랫폼", 
                "기업 적용 성공 사례", 
                "구현 시 주요 도전과제"
            ],
            "image": None
        },
        # 슬라이드 34: 미래 발전 방향
        {
            "title": "미래 발전 방향",
            "bullets": [
                "RAG와 에이전트의 진화 전망", 
                "새로운 연구 분야와 기회", 
                "산업별 적용 가능성", 
                "윤리적 고려사항"
            ],
            "image": None
        },
        # 슬라이드 35: 결론
        {
            "title": "결론",
            "bullets": [
                "핵심 메시지 요약", 
                "RAG와 에이전트의 실용적 중요성", 
                "구현 시작을 위한 조언", 
                "미래 연구 및 발전 방향"
            ],
            "image": None
        },
        # 슬라이드 36: 참고 자료
        {
            "title": "참고 자료",
            "bullets": [
                "추천 논문, 책, 블로그", 
                "오픈소스 프로젝트 및 라이브러리", 
                "학습 리소스", 
                "연락처 및 질문"
            ],
            "image": None
        },
        # 슬라이드 37: Q&A
        {
            "title": "Q&A",
            "bullets": [
                "질문 및 토론 안내", 
                "추가 정보 요청 방법", 
                "피드백 수집"
            ],
            "image": None
        }
    ]
    
    # 슬라이드 생성 함수
    def add_content_slide(title, bullets, image_filename=None):
        # 제목과 내용이 있는 레이아웃 선택
        content_slide_layout = prs.slide_layouts[1]  # 제목 및 내용 슬라이드 레이아웃
        slide = prs.slides.add_slide(content_slide_layout)
        
        # 제목 설정
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # 폰트 설정
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(36)
                run.font.bold = True
        
        # 내용 설정
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()  # 기존 내용 제거
        
        # 불릿 포인트 추가
        for bullet_text in bullets:
            p = text_frame.add_paragraph()
            p.text = bullet_text
            p.level = 0  # 첫 번째 레벨 (들여쓰기 없음)
            p.font.name = font_name
            p.font.size = Pt(24)  # 이미지가 있을 경우 글자 크기를 조정
        
        # 이미지 추가 (있는 경우)
        if image_filename:
            img_path = os.path.join(IMAGE_DIR, image_filename)
            if os.path.exists(img_path):
                # 이미지의 위치와 크기 설정
                left = Inches(6.0)  # 오른쪽에 배치
                top = Inches(2.5)   # 약간 아래로
                width = Inches(4.0)  # 너비
                
                # 이미지 추가
                slide.shapes.add_picture(img_path, left, top, width=width)
    
    # 모든 슬라이드 콘텐츠 추가
    for content in slide_contents:
        add_content_slide(content["title"], content["bullets"], content["image"])
    
    # 프레젠테이션 저장
    output_path = "/Users/macbook/Documents/RAG와_에이전트_프레젠테이션.pptx"
    prs.save(output_path)
    print(f"프레젠테이션이 저장되었습니다: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_presentation() 